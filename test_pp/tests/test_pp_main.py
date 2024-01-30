#kevin fink
#test_pp_main.py
#Jan 30 2024
#kevin@shorecode.org

import unittest
from unittest import mock
import pp_main
from pp_main import query_chatgpt
from pptx import Presentation
from pptx.slide import SlideLayout

class PowerPointLayoutTest(unittest.TestCase):

    def setUp(self):
        self.presentation_filepath = 'data/GIS_Roca_Advanced_Ceramic_Market_Report - Copy.pptx'
        self.presentation = Presentation(self.presentation_filepath)

    def test_alter_section(self):
        new_content = 'Testing'
        pp_main.alter_section_content(self.presentation, 3, new_content)
        self.assertEqual(new_content, self.presentation.slides[3].shapes[4].text_frame.text)
        
    def test_query_chatgpt(self):
        prompt = "Hello, how are you?"
        assistant_profile = "My name is John, I'm 25 years old."
        
        with mock.patch('openai.chat.completions.create') as mock_create:
            mock_response = mock.Mock()
            mock_message = mock.Mock()
            mock_choice = mock.Mock()
            
            mock_response.choices = [mock_choice]
            mock_choice.message.content.strip.return_value = "I'm fine, thank you."
            
            mock_create.return_value = mock_response
            
            result = query_chatgpt(prompt, assistant_profile)
            
            mock_create.assert_called_once_with(
                model='gpt-4',
                messages=[{'role': 'system', 'content': assistant_profile},
                          {'role': 'user', 'content': prompt}],
                n=1,
                stop=None,
                temperature=0.7
            )
            assert result == "I'm fine, thank you."


if __name__ == '__main__':
    unittest.main()