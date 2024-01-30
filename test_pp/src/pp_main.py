#kevin fink
#pp_main.py
#Jan 30 2024
#kevin@shorecode.org

import datetime
import os
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches
import openai
import pp_logging

logger = pp_logging.set_logging('test_pp', 'pp_main.log')

def alter_section_content(presentation: Presentation, section_index: int, new_content: str):
    '''
    This function alters the content of a section within a presentation.

    param presentation (Presentation): the PowerPoint presentation object 
    param section_index (int): the index of the section to be altered 
    param new_content (str): the new content to be set for the section 
    return: None
    '''
    logger.info(f'Changing slide content for slide # {section_index+1}')
    slides = presentation.slides
    logger.info(new_content)
    for slide in slides:
        # Determines if the slide matches the provided section_index
        if slide.slide_id == presentation.slides[section_index].slide_id:
            if slide.shapes[4].has_text_frame:
                slide.shapes[4].text_frame.clear()  # Clear existing content
                slide.shapes[4].text_frame.text = new_content # Set new content)
                slide.shapes[4].text_frame.fit_text
                # Iterate through all the paragraphs in the text box
                for paragraph in slide.shapes[4].text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Change the color of the text to new_color
                        run.font.color.rgb = RGBColor(255,0,0)                
                logger.info(f'Slide # {section_index+1} succesfully changed')
            slide.shapes[5].text_frame.clear()


def query_chatgpt(prompt: str, assistant_profile: str) -> str:
    '''
    The query_chatgpt function takes a prompt and an assistant_profile as inputs and uses OpenAI's Chat API to generate a response.
    param prompt (str): The user's input prompt.
    param assistant_profile (str): Information about the assistant (e.g., name, age, occupation). 
    return (str): The response generated by the Chat API.
    '''
    model = 'gpt-4'
    logger.info(f'Sending request to {model}')
    response = openai.chat.completions.create(
    model=model,
    messages=[{'role': 'system', 'content': assistant_profile},
              {'role': 'user', 'content': prompt}],
    n=1,
    stop=None,
    temperature=0.7
    )
    return response.choices[0].message.content.strip()

# Set up OpenAI API credentials
openai.api_key = os.getenv('OPENAI_API_KEY')

# Load the PowerPoint presentation
presentation_filepath = 'data/GIS_Roca_Advanced_Ceramic_Market_Report - Copy.pptx'
logger.info('Loading the powerpoint presentation')
presentation = Presentation(presentation_filepath)

# Get response from ChatGPT
prompt = 'Write a market definition for the Medical Imaging Market in 230 words or less. Then write a \
summary of common applications for the Medical Imaging market in 150 words or less using \
a bulleted list with 3 points. Do not include a title for this definition, only the content.'
assistant_profile = 'You are a professional consultant that is an expert in market \
research. Your target audience is business executives. Be concise and only include \
information that is factual.'
new_content = query_chatgpt(prompt, assistant_profile)
logger.info('OpenAI response received')

# Alter the content of a specific section (e.g., section at index 2) with new content
alter_section_content(presentation, 3, new_content)

# Get the current date
current_date = datetime.datetime.now()

# Format the date as "DD_MM_YY"
formatted_date = current_date.strftime("%d_%m_%y")

# Get the current time
current_time = datetime.datetime.now().strftime("%H%M")

# Save the modified presentation
logger.info('Saving new PowerPoint file')
presentation.save(f'{presentation_filepath}_modified_{formatted_date}_{current_time}.pptx')